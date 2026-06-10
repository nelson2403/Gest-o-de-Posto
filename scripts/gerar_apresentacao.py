# -*- coding: utf-8 -*-
"""Gera a apresentação (treinamento interno) do Sistema de Gestão Pedra do Pombal."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Cores (identidade do sistema) ──
VINHO   = RGBColor(0x8B, 0x1A, 0x14)
LARANJA = RGBColor(0xF9, 0x73, 0x16)
ESCURO  = RGBColor(0x1F, 0x29, 0x37)
CINZA   = RGBColor(0x6B, 0x72, 0x80)
CLARO   = RGBColor(0xF3, 0xF4, 0xF6)
BRANCO  = RGBColor(0xFF, 0xFF, 0xFF)
FONTE   = "Calibri"

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO = os.path.join(BASE, "public", "logo.png")
SAIDA = os.path.join(BASE, "Apresentacao_Sistema_Pedra_do_Pombal.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _no_line(shape):
    shape.line.fill.background()


def retangulo(slide, left, top, width, height, cor):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sp.fill.solid()
    sp.fill.fore_color.rgb = cor
    _no_line(sp)
    sp.shadow.inherit = False
    return sp


def caixa_texto(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(run, texto, tam, cor, bold=False, italic=False):
    run.text = texto
    run.font.size = Pt(tam)
    run.font.color.rgb = cor
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONTE


def rodape(slide, num):
    tb, tf = caixa_texto(slide, Inches(0.4), Inches(7.05), Inches(9), Inches(0.35))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Sistema de Gestão • Pedra do Pombal", 9, CINZA)
    tb2, tf2 = caixa_texto(slide, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.35))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    set_run(p2.add_run(), str(num), 9, CINZA)


# ── Slide de capa ──
def slide_capa():
    s = prs.slides.add_slide(BLANK)
    retangulo(s, 0, 0, SW, SH, VINHO)
    retangulo(s, 0, Inches(6.9), SW, Inches(0.6), LARANJA)
    # logo
    if os.path.exists(LOGO):
        s.shapes.add_picture(LOGO, Inches(0.7), Inches(0.6), height=Inches(1.5))
    tb, tf = caixa_texto(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.4))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Sistema de Gestão de Postos", 44, BRANCO, bold=True)
    p2 = tf.add_paragraph()
    set_run(p2.add_run(), "Pedra do Pombal", 30, RGBColor(0xFF, 0xC9, 0xB8), bold=True)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(18)
    set_run(p3.add_run(), "Treinamento Interno — Visão dos Módulos", 18, RGBColor(0xF0, 0xD8, 0xD4))
    return s


# ── Slide de conteúdo ──
def slide_conteudo(titulo, subtitulo, bullets, num):
    s = prs.slides.add_slide(BLANK)
    # barra superior
    retangulo(s, 0, 0, SW, Inches(1.25), VINHO)
    retangulo(s, 0, Inches(1.25), SW, Inches(0.09), LARANJA)
    # título
    tb, tf = caixa_texto(s, Inches(0.5), Inches(0.18), Inches(12.3), Inches(1.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    set_run(p.add_run(), titulo, 30, BRANCO, bold=True)
    if subtitulo:
        p2 = tf.add_paragraph()
        set_run(p2.add_run(), subtitulo, 14, RGBColor(0xF0, 0xD8, 0xD4))
    # bullets
    tb2, tf2 = caixa_texto(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(5.1))
    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            texto, nivel = item
        else:
            texto, nivel = item, 0
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(8)
        if nivel == 0:
            run_b = p.add_run()
            set_run(run_b, "●  ", 16, LARANJA, bold=True)
            set_run(p.add_run(), texto, 17, ESCURO, bold=True)
        else:
            p.level = 1
            run_b = p.add_run()
            set_run(run_b, "–  ", 14, CINZA)
            set_run(p.add_run(), texto, 14, RGBColor(0x37, 0x41, 0x51))
    rodape(s, num)
    return s


# ── Slide de encerramento ──
def slide_fim():
    s = prs.slides.add_slide(BLANK)
    retangulo(s, 0, 0, SW, SH, VINHO)
    retangulo(s, 0, 0, SW, Inches(0.6), LARANJA)
    tb, tf = caixa_texto(s, Inches(0.9), Inches(2.6), Inches(11.5), Inches(2.4))
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Tudo num só lugar.", 40, BRANCO, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(14)
    set_run(p2.add_run(),
            "Controle financeiro, fiscal, estoque, caixa e conciliação — integrados ao AUTOSYSTEM.",
            18, RGBColor(0xF0, 0xD8, 0xD4))
    p3 = tf.add_paragraph()
    p3.space_before = Pt(24)
    set_run(p3.add_run(), "Dúvidas no uso? Fale com o setor financeiro/TI.", 15,
            RGBColor(0xFF, 0xC9, 0xB8))
    return s


# ─────────────────────────────────────────────────────────────────────────
slide_capa()

slide_conteudo(
    "Visão Geral do Sistema",
    "O painel central de gestão dos postos",
    [
        "Reúne, em um só lugar, a gestão financeira, fiscal, de estoque e de caixa dos postos.",
        "Integrado ao AUTOSYSTEM — puxa vendas, movimentos e títulos automaticamente.",
        "Acesso pelo navegador, com menu por áreas no topo:",
        ("Dashboard, Financeiro, Fiscal, Compras, Comissionamento, Controle Geral e Analítico.", 1),
        "Cada usuário vê apenas o que seu perfil de acesso permite.",
    ],
    2,
)

slide_conteudo(
    "Financeiro",
    "Contas, caixas e dinheiro sob controle",
    [
        "Contas a Receber: acompanhamento de recebíveis por posto e forma de pagamento.",
        "Contas a Pagar: conferência diária, despesas fixas, títulos e boletos/solicitações.",
        "Controle de Caixas e Controle de Dinheiro: saldos e movimentações por posto.",
        "Fechamento de Caixa Eletrônico: consulta dos fechamentos enviados pelos frentistas.",
    ],
    3,
)

slide_conteudo(
    "Conciliação Bancária & Divergências",
    "Conferência automática extrato × sistema",
    [
        "Geração de tarefas de conciliação e painel do Extrato Bancário.",
        "O sistema compara o movimento do extrato com o do AUTOSYSTEM e aponta divergências.",
        "Conciliadores veem apenas as divergências dos seus postos; o master vê todas.",
        "Ao entrar no sistema, o recálculo das divergências é feito automaticamente.",
        ("Clicar na divergência leva direto à tarefa para resolver.", 1),
    ],
    4,
)

slide_conteudo(
    "Fiscal",
    "Notas, manifestos e boletos",
    [
        "Painel Fiscal: visão das notas e situação fiscal por posto.",
        "Tarefas Fiscal: reconhecimento de NF, boletos, romaneios e uso/consumo.",
        "Manifestação de notas e sincronização automática com o AUTOSYSTEM.",
        "Envio de boletos para o Contas a Pagar de forma integrada.",
    ],
    5,
)

slide_conteudo(
    "Compras & Estoque",
    "Do pedido à contagem",
    [
        "Estoque de combustíveis e conveniência, com contagem por pista/produto.",
        "Uso e Consumo: controle das despesas marcadas para consumo interno.",
        "Sugestão de Pedido: ajuda a decidir o que comprar com base no giro.",
        "Cadastro de Fornecedores e rotina de visitas.",
    ],
    6,
)

slide_conteudo(
    "Fechamento de Caixa (PDV)",
    "Tela do frentista, separada do sistema principal",
    [
        "O frentista acessa /pdv e entra com código de operador + PIN.",
        "Preenche os valores que tem em caixa; o sistema revela os valores do AUTOSYSTEM e mostra a diferença.",
        "Assina na tela e o comprovante é impresso na impressora térmica (via QZ Tray).",
        "Regras de segurança:",
        ("Só é permitido fechar o dia atual.", 1),
        ("Apenas um fechamento por frentista por dia.", 1),
    ],
    7,
)

slide_conteudo(
    "Comissionamento & Analítico",
    "Metas, comissões e análise de vendas",
    [
        "Comissionamento: esquemas, regras, metas, simulação e relatórios (perfil master).",
        "Analítico e Análise de Vendas: acompanhamento de desempenho por posto e período.",
        "Dashboards para apoiar decisões do dia a dia.",
    ],
    8,
)

slide_conteudo(
    "Acessos & Segurança",
    "Cada um vê só o que precisa",
    [
        "Perfis de acesso definem o que cada usuário enxerga e pode fazer.",
        "Papéis como master, administradores, conciliador e operador de caixa.",
        "Conciliador acessa apenas os postos atribuídos a ele.",
        "Central de Acessos: portais, senhas, câmeras e servidores organizados.",
    ],
    9,
)

slide_fim()

prs.save(SAIDA)
print("OK:", SAIDA)
